# -*- coding: utf-8 -*-
"""
文档脱敏工具 v4.1
依赖：pip install python-docx python-pptx openpyxl tkinterdnd2
"""

import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import uuid, json, os, re, threading

try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    DND_OK = True
except ImportError:
    DND_OK = False

try:
    from docx import Document as DocxDoc
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

MAPPING_SUFFIX = "_mask_mapping.json"
SUPPORTED_EXT  = {".docx", ".pptx", ".xlsx", ".txt"}

# ── 中英文语言包 ──
LANG = {
    "zh": {
        "title":            "文档脱敏",
        "subtitle":         "本地处理  ·  数据不上传",
        "badge_empty":      "映射：空",
        "badge_n":          "映射：{n} 条",
        "sensitive_title":  "敏感词",
        "sensitive_hint":   "多个词请用逗号分隔（支持中文逗号）",
        "sensitive_ph":     "例：中国，1，2，3，0，政府",
        "enc_zone":         "加密区",
        "dec_zone":         "解密区",
        "start_enc":        "开始加密",
        "start_dec":        "开始解密",
        "clear_map":        "清空映射",
        "copy_prompt":      "复制AI提示词",
        "del_json_opt":     "解密成功后自动删除映射文件（JSON）",
        "ready":            "就绪",
        "lang_btn":         "EN",
        "add_file":         "添加文件",
        "add_folder":       "添加文件夹",
        "clear":            "清空",
        "load_json":        "载入JSON映射",
        "hint_enc":         "将文件拖入此处，或点击下方按钮选择",
        "hint_dec":         "将文件或 JSON 映射拖入此处",
        "footer":           "与 AI 对话前，请先复制提示词，并粘贴到AI对话框。",
        "dep_title":        "缺少依赖",
        "dep_body":         "以下库未安装，对应功能不可用：\n\n{libs}\n\n安装：\npip install python-docx python-pptx openpyxl tkinterdnd2",
        "warn_no_file_enc": "请先在加密区添加文件。",
        "warn_no_file_dec": "请先在解密区添加文件。",
        "warn_no_words":    "请填写敏感词。",
        "warn_no_map_title":"未找到映射文件",
        "warn_no_map_body": "未在文件目录找到映射文件。\n\n点击「确定」手动选择 JSON 映射文件。",
        "processing":       "处理中...",
        "enc_done":         "加密",
        "dec_done":         "解密",
        "success":          "成功",
        "fail":             "失败",
        "map_loaded":       "映射已载入，共 {n} 条",
        "map_load_ok_title":"载入成功",
        "map_load_ok_body": "映射文件已加载！共 {n} 条记录。\n现在可将文件拖入解密区执行解密。",
        "map_cleared":      "映射已清空",
        "map_clear_title":  "已清空",
        "map_clear_body":   "映射关系已清空，可开始新任务。",
        "copied_title":     "已复制到剪贴板",
        "copied_body":      "使用方法：\n\n① 在 AI 对话框最前面粘贴此提示词\n② 再粘贴 / 上传脱敏后的文件内容\n③ AI 将原样保留所有 [MASK_XXXXXXXX] 标记\n\n收到 AI 回传内容后，保存为文件，\n拖入解密区一键还原。",
        "auto_load_map":    "自动载入映射：{name}",
        "map_deleted":      "映射文件已删除：{names}",
        "suggest_prompt":   "\n建议：发给 AI 前，点击「复制AI提示词」，\n将提示词粘贴在对话最前面，AI 会保留所有占位符不变。\n本次映射已保存在内存，收到回传文件后直接拖入解密区即可。",
        "all_ok":           " 全部成功",
        "partial_fail":     " 完成（含失败项）",
        "ok_n":             "成功 {n} 个：",
        "fail_n":           "失败 {n} 个：",
        "pick_file_title":  "选择文件",
        "pick_file_types":  [("支持的文件", "*.docx *.pptx *.xlsx *.txt"), ("所有文件", "*.*")],
        "pick_dir_title":   "选择文件夹",
        "pick_json_title":  "选择 JSON 映射文件",
        "pick_map_title":   "选择映射文件",
        "load_fail":        "载入失败",
        "error":            "错误",
        "tip":              "提示",
    },
    "en": {
        "title":            "Doc Mask",
        "subtitle":         "Local Only  ·  No Upload",
        "badge_empty":      "Map: empty",
        "badge_n":          "Map: {n}",
        "sensitive_title":  "Sensitive Words",
        "sensitive_hint":   "Separate multiple words with commas",
        "sensitive_ph":     "e.g. China, 1, 2, 3, gov",
        "enc_zone":         "Encrypt Zone",
        "dec_zone":         "Decrypt Zone",
        "start_enc":        "Start Encrypt",
        "start_dec":        "Start Decrypt",
        "clear_map":        "Clear Map",
        "copy_prompt":      "Copy AI Prompt",
        "del_json_opt":     "Auto-delete mapping file (JSON) after successful decryption",
        "ready":            "Ready",
        "lang_btn":         "中",
        "add_file":         "Add File",
        "add_folder":       "Add Folder",
        "clear":            "Clear",
        "load_json":        "Load JSON Map",
        "hint_enc":         "Drag files here or click button below",
        "hint_dec":         "Drag files or JSON map here",
        "footer":           "Before chatting with AI, copy and paste the prompt first.",
        "dep_title":        "Missing Dependencies",
        "dep_body":         "The following libraries are not installed:\n\n{libs}\n\nInstall:\npip install python-docx python-pptx openpyxl tkinterdnd2",
        "warn_no_file_enc": "Please add files to the Encrypt Zone first.",
        "warn_no_file_dec": "Please add files to the Decrypt Zone first.",
        "warn_no_words":    "Please enter sensitive words.",
        "warn_no_map_title":"Mapping File Not Found",
        "warn_no_map_body": "No mapping file found in file directory.\n\nClick OK to manually select a JSON mapping file.",
        "processing":       "Processing...",
        "enc_done":         "Encryption",
        "dec_done":         "Decryption",
        "success":          "succeeded",
        "fail":             "failed",
        "map_loaded":       "Map loaded: {n} entries",
        "map_load_ok_title":"Load Successful",
        "map_load_ok_body": "Mapping loaded! {n} entries.\nNow drag files to Decrypt Zone.",
        "map_cleared":      "Map cleared",
        "map_clear_title":  "Cleared",
        "map_clear_body":   "Mapping cleared. Ready for a new task.",
        "copied_title":     "Copied to Clipboard",
        "copied_body":      "How to use:\n\n① Paste this prompt at the start of your AI chat\n② Then paste / upload the masked file content\n③ AI will keep all [MASK_XXXXXXXX] tags intact\n\nAfter receiving the AI reply, save as a file and drag into Decrypt Zone.",
        "auto_load_map":    "Auto-loaded map: {name}",
        "map_deleted":      "Map file deleted: {names}",
        "suggest_prompt":   "\nTip: Before sending to AI, click 'Copy AI Prompt' and paste it at the start of the conversation. The map is saved in memory — drag the returned file into Decrypt Zone to restore.",
        "all_ok":           " All Succeeded",
        "partial_fail":     " Completed (with failures)",
        "ok_n":             "Succeeded {n}:",
        "fail_n":           "Failed {n}:",
        "pick_file_title":  "Select Files",
        "pick_file_types":  [("Supported Files", "*.docx *.pptx *.xlsx *.txt"), ("All Files", "*.*")],
        "pick_dir_title":   "Select Folder",
        "pick_json_title":  "Select JSON Mapping File",
        "pick_map_title":   "Select Mapping File",
        "load_fail":        "Load Failed",
        "error":            "Error",
        "tip":              "Notice",
    },
}

AI_PROMPT = (
    "【重要指令】本文档中形如 [MASK_XXXXXXXX] 的标记是敏感信息占位符，"
    "由脱敏系统自动生成。处理本文档时必须遵守：\n"
    "1. 严禁修改、翻译、解释或删除任何 [MASK_XXXXXXXX] 标记；\n"
    "2. 将每个标记视为不可分割的词语，原样保留在输出中；\n"
    "3. 仅对标记以外的正常文字执行你的任务。\n"
    "违反上述规则将导致解密失败。"
)

# ── 配色（Apple 风格）──
BG      = "#F2F2F7"
CARD    = "#FFFFFF"
BORDER  = "#D1D1D6"
T1      = "#1C1C1E"
T2      = "#6C6C70"
ORANGE  = "#FF9500"
BLUE    = "#007AFF"
GREEN   = "#34C759"
PURPLE  = "#5856D6"
RED     = "#FF3B30"
FONT    = "Microsoft YaHei"


# ==============================================================================
# 核心引擎
# ==============================================================================

class Scrubber:

    def __init__(self):
        self.mapping  = {}   # word -> mask
        self.reverse  = {}   # mask -> word
        self.json_paths = [] # 已保存的 JSON 路径列表

    def build(self, words):
        for w in words:
            w = w.strip()
            if w and w not in self.mapping:
                m = "[MASK_" + uuid.uuid4().hex[:8].upper() + "]"
                self.mapping[w] = m
                self.reverse[m] = w

    def _enc(self, text):
        """
        加密时跳过已有掩码段，
        且使用单次正则替换，避免二次污染掩码内容。
        """

        # 匹配已有 MASK
        mask_pat = re.compile(r'\[MASK_[0-9A-F]{8}\]')

        # 敏感词为空直接返回
        if not self.mapping:
            return text

        # 构造敏感词正则（按长度降序避免短词抢长词）
        words = sorted(self.mapping.keys(), key=len, reverse=True)
        word_pat = re.compile("|".join(re.escape(w) for w in words))

        parts = mask_pat.split(text)
        masks = mask_pat.findall(text)

        safe_parts = []

        for part in parts:
            # 在纯文本片段中做一次性替换
            def repl(match):
                return self.mapping.get(match.group(0), match.group(0))

            new_part = word_pat.sub(repl, part)
            safe_parts.append(new_part)

        # 拼接回去
        out = safe_parts[0]
        for mk, p in zip(masks, safe_parts[1:]):
            out += mk + p

        return out

    def _dec(self, text):
        pat = re.compile(r'\[MASK_[0-9A-F]{8}\]')
        return pat.sub(lambda m: self.reverse.get(m.group(0), m.group(0)), text)

    def save(self, out_path):
        d    = os.path.dirname(out_path) or "."
        base = os.path.splitext(os.path.basename(out_path))[0]
        mp   = os.path.join(d, base + MAPPING_SUFFIX)
        with open(mp, "w", encoding="utf-8") as f:
            json.dump({"mapping": self.mapping, "reverse": self.reverse},
                      f, ensure_ascii=False, indent=2)
        self.json_paths.append(mp)
        return mp

    def load(self, path):
        with open(path, "r", encoding="utf-8") as f:
            d = json.load(f)
        self.mapping.update(d.get("mapping", {}))
        self.reverse.update(d.get("reverse", d.get("reverse_mapping", {})))

    def auto_load(self, doc_path):
        d    = os.path.dirname(doc_path) or "."
        base = os.path.splitext(os.path.basename(doc_path))[0]
        mp   = os.path.join(d, base + MAPPING_SUFFIX)
        if os.path.exists(mp):
            self.load(mp)
            return True, mp
        return False, mp

    # ── 文档处理 ──

    def _handle_docx(self, src, dst, fn):
        if not DOCX_OK:
            raise ImportError("请安装 python-docx")
        doc = DocxDoc(src)
        def fix(p):
            full = p.text
            new  = fn(full)
            if full == new: return
            for r in p.runs: r.text = fn(r.text)
            if p.text != new:
                for i, r in enumerate(p.runs):
                    r.text = new if i == 0 else ""
        for p in doc.paragraphs: fix(p)
        for t in doc.tables:
            for row in t.rows:
                for c in row.cells:
                    for p in c.paragraphs: fix(p)
        doc.save(dst)

    def _handle_pptx(self, src, dst, fn):
        if not PPTX_OK:
            raise ImportError("请安装 python-pptx")
        prs = Presentation(src)
        def fix_tf(tf):
            for p in tf.paragraphs:
                full = "".join(r.text for r in p.runs)
                new  = fn(full)
                if full == new: continue
                for r in p.runs: r.text = fn(r.text)
                if "".join(r.text for r in p.runs) != new:
                    for i, r in enumerate(p.runs):
                        r.text = new if i == 0 else ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame: fix_tf(shape.text_frame)
                if shape.has_table:
                    for row in shape.table.rows:
                        for c in row.cells: fix_tf(c.text_frame)
        prs.save(dst)

    def _handle_xlsx(self, src, dst, fn):
        if not XLSX_OK:
            raise ImportError("请安装 openpyxl")
        wb = openpyxl.load_workbook(src)
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for c in row:
                    if isinstance(c.value, str):
                        c.value = fn(c.value)
        wb.save(dst)

    def _handle_txt(self, src, dst, fn):
        content = ""
        for enc in ("utf-8", "utf-8-sig", "gbk"):
            try:
                with open(src, "r", encoding=enc) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        with open(dst, "w", encoding="utf-8") as f:
            f.write(fn(content))

    def process(self, src, words, mode):
        if not os.path.exists(src):
            raise FileNotFoundError("文件不存在：" + src)
        ext  = os.path.splitext(src)[1].lower()
        d    = os.path.dirname(src) or "."
        base = os.path.splitext(os.path.basename(src))[0]
        dst  = os.path.join(d, base + ("_加密" if mode == "enc" else "_解密") + ext)

        fn = self._enc if mode == "enc" else self._dec

        H = {".docx": self._handle_docx, ".pptx": self._handle_pptx,
             ".xlsx": self._handle_xlsx,  ".txt":  self._handle_txt}
        if ext in H:
            H[ext](src, dst, fn)
        elif ext in {".doc", ".ppt", ".xls"}:
            raise ValueError(ext + " 为旧版格式，请另存为新格式后处理")
        else:
            raise ValueError("不支持的格式：" + ext)

        mp = self.save(dst) if mode == "enc" else None
        return dst, mp

    def batch(self, files, words, mode, cb=None):
        if mode == "enc":
            self.build(words)
        results = []
        for i, fp in enumerate(files):
            fname = os.path.basename(fp)
            try:
                out, mp = self.process(fp, words, mode)
                results.append((fp, out, None, mp))
                if cb: cb(i+1, len(files), fname, "完成")
            except Exception as e:
                results.append((fp, None, str(e), None))
                if cb: cb(i+1, len(files), fname, "失败")
        return results


# ==============================================================================
# 文件列表组件（纯 pack 布局）
# ==============================================================================

class FileList(tk.Frame):
    """
    文件拖放列表，严格使用 pack 布局，避免与外部 grid 冲突
    """

    def __init__(self, master, accent, mode, scrubber,
                 on_json=None, lang_var=None, **kw):
        super().__init__(master, bg=CARD, **kw)
        self.accent   = accent
        self.mode     = mode        # "enc" | "dec"
        self.scrubber = scrubber
        self.on_json  = on_json
        self.lang_var = lang_var    # StringVar "zh"/"en"
        self.files    = []
        self._build()
        if DND_OK:
            self._bind_dnd()

    def _t(self, key):
        lang = self.lang_var.get() if self.lang_var else "zh"
        return LANG[lang][key]

    def _build(self):
        # 拖放提示
        hint_bg = self._tint(self.accent, 0.92)
        hint_key = "hint_enc" if self.mode == "enc" else "hint_dec"
        self.hint = tk.Label(
            self, bg=hint_bg, fg=self.accent, cursor="hand2",
            font=(FONT, 9),
            text=self._t(hint_key)
        )
        self.hint.pack(fill="x", ipady=10)
        self.hint.bind("<Button-1>", lambda e: self._pick())

        # 列表区（用 Frame 包裹 Listbox + Scrollbar，全用 pack）
        list_frame = tk.Frame(self, bg=CARD)
        list_frame.pack(fill="both", expand=True, pady=4)

        sb = tk.Scrollbar(list_frame)
        sb.pack(side="right", fill="y")

        self.lb = tk.Listbox(
            list_frame, font=(FONT, 9),
            yscrollcommand=sb.set,
            selectmode=tk.EXTENDED,
            relief="flat", bd=0,
            bg="#F8F8FA", fg=T1,
            selectbackground=self.accent,
            selectforeground="white",
            activestyle="none"
        )
        self.lb.pack(side="left", fill="both", expand=True, padx=(4, 0))
        sb.config(command=self.lb.yview)

        self.lb.bind("<Delete>",    self._del)
        self.lb.bind("<BackSpace>", self._del)
        self.lb.bind("<Button-3>",  self._ctx)

        # 按钮行（全用 pack）
        self.br = tk.Frame(self, bg=CARD)
        self.br.pack(fill="x", pady=(0, 4))

        self.btn_add_file   = self._mkbtn(self.br, self._t("add_file"),   self.accent,  self._pick)
        self.btn_add_folder = self._mkbtn(self.br, self._t("add_folder"), self.accent,  self._pick_dir)
        self.btn_clear_list = self._mkbtn(self.br, self._t("clear"),      "#AEAEB2",    self._clear)
        self.btn_add_file.pack(side="left", padx=(4, 3))
        self.btn_add_folder.pack(side="left", padx=(0, 3))
        self.btn_clear_list.pack(side="left")

        self.btn_load_json = None
        if self.mode == "dec":
            self.btn_load_json = self._mkbtn(self.br, self._t("load_json"), PURPLE, self._pick_json)
            self.btn_load_json.pack(side="right", padx=4)

        self.cnt = tk.Label(self.br, text="0 " + ("个" if (self.lang_var and self.lang_var.get()=="zh") else "files"),
                            font=(FONT, 8), bg=CARD, fg=T2)
        self.cnt.pack(side="right", padx=6)

    def update_lang(self):
        """切换语言时刷新所有文字"""
        hint_key = "hint_enc" if self.mode == "enc" else "hint_dec"
        self.hint.config(text=self._t(hint_key))
        # 重建按钮（Canvas 无法直接改文字，销毁重建）
        for w in self.br.winfo_children():
            w.destroy()
        self.btn_add_file   = self._mkbtn(self.br, self._t("add_file"),   self.accent,  self._pick)
        self.btn_add_folder = self._mkbtn(self.br, self._t("add_folder"), self.accent,  self._pick_dir)
        self.btn_clear_list = self._mkbtn(self.br, self._t("clear"),      "#AEAEB2",    self._clear)
        self.btn_add_file.pack(side="left", padx=(4, 3))
        self.btn_add_folder.pack(side="left", padx=(0, 3))
        self.btn_clear_list.pack(side="left")
        if self.mode == "dec":
            self.btn_load_json = self._mkbtn(self.br, self._t("load_json"), PURPLE, self._pick_json)
            self.btn_load_json.pack(side="right", padx=4)
        n = len(self.files)
        suffix = "个" if self._t("lang_btn") == "EN" else "files"
        self.cnt = tk.Label(self.br, text=str(n) + " " + suffix,
                            font=(FONT, 8), bg=CARD, fg=T2)
        self.cnt.pack(side="right", padx=6)

    def _mkbtn(self, p, txt, col, cmd):
        import tkinter.font as tkfont
        f = tkfont.Font(family=FONT, size=8)
        tw = f.measure(txt)
        th = f.metrics("linespace")
        padx, pady, radius = 8, 4, 8
        w = tw + padx * 2
        h = th + pady * 2
        hover = self._tint(col, 0.82)

        cv = tk.Canvas(p, width=w, height=h,
                       bg=p["bg"], highlightthickness=0, cursor="hand2")

        def _draw(c):
            cv.delete("all")
            r = radius
            cv.create_arc(0, 0, 2*r, 2*r, start=90, extent=90, fill=c, outline=c)
            cv.create_arc(w-2*r, 0, w, 2*r, start=0, extent=90, fill=c, outline=c)
            cv.create_arc(0, h-2*r, 2*r, h, start=180, extent=90, fill=c, outline=c)
            cv.create_arc(w-2*r, h-2*r, w, h, start=270, extent=90, fill=c, outline=c)
            cv.create_rectangle(r, 0, w-r, h, fill=c, outline=c)
            cv.create_rectangle(0, r, w, h-r, fill=c, outline=c)
            cv.create_text(w//2, h//2, text=txt, fill="white",
                           font=(FONT, 8))

        _draw(col)
        cv.bind("<Button-1>", lambda e: cmd())
        cv.bind("<Enter>",    lambda e: _draw(hover))
        cv.bind("<Leave>",    lambda e: _draw(col))
        return cv

    @staticmethod
    def _tint(hex_c, factor):
        h = hex_c.lstrip("#")
        r = int(int(h[0:2], 16) * factor + 255 * (1-factor))
        g = int(int(h[2:4], 16) * factor + 255 * (1-factor))
        b = int(int(h[4:6], 16) * factor + 255 * (1-factor))
        return "#{:02X}{:02X}{:02X}".format(
            min(r, 255), min(g, 255), min(b, 255))

    # ── DND ──
    def _bind_dnd(self):
        for w in (self, self.hint, self.lb):
            w.drop_target_register(DND_FILES)
            w.dnd_bind("<<Drop>>", self._drop)

    def _drop(self, ev):
        raw   = ev.data
        paths = re.findall(r'\{([^}]+)\}|(\S+)', raw)
        flat  = [p[0] or p[1] for p in paths]
        for p in flat:
            p = p.strip().strip('"')
            if not p:
                continue
            if p.lower().endswith(".json") and self.mode == "dec":
                self._load_json(p)
            elif os.path.isdir(p):
                self._scan_dir(p)
            else:
                self._add(p)
        self._refresh()

    # ── 文件操作 ──
    def _pick(self):
        ps = filedialog.askopenfilenames(
            title=self._t("pick_file_title"),
            filetypes=self._t("pick_file_types")
        )
        for p in ps: self._add(p)
        self._refresh()

    def _pick_dir(self):
        d = filedialog.askdirectory(title=self._t("pick_dir_title"))
        if d:
            self._scan_dir(d)
            self._refresh()

    def _pick_json(self):
        p = filedialog.askopenfilename(
            title=self._t("pick_json_title"),
            filetypes=[("JSON", "*.json"), ("All Files" if self._t("lang_btn")=="中" else "所有文件", "*.*")]
        )
        if p: self._load_json(p)

    def _load_json(self, path):
        try:
            self.scrubber.load(path)
            n = len(self.scrubber.mapping)
            if self.on_json:
                self.on_json(n, path)
        except Exception as e:
            messagebox.showerror(self._t("load_fail"), str(e))

    def _scan_dir(self, d):
        for fn in os.listdir(d):
            if os.path.splitext(fn)[1].lower() in SUPPORTED_EXT:
                self._add(os.path.join(d, fn))

    def _add(self, path):
        if not path: return
        if os.path.splitext(path)[1].lower() not in SUPPORTED_EXT: return
        if path in self.files: return
        self.files.append(path)
        self.lb.insert(tk.END, "  " + os.path.basename(path))

    def _del(self, ev=None):
        for i in list(self.lb.curselection())[::-1]:
            self.lb.delete(i)
            self.files.pop(i)
        self._refresh()

    def _ctx(self, ev):
        m = tk.Menu(self, tearoff=0)
        m.add_command(label="删除选中" if self._t("lang_btn")=="EN" else "Remove Selected", command=self._del)
        m.add_command(label="清空列表" if self._t("lang_btn")=="EN" else "Clear List", command=self._clear)
        m.tk_popup(ev.x_root, ev.y_root)

    def _clear(self):
        self.files.clear()
        self.lb.delete(0, tk.END)
        self._refresh()

    def _refresh(self):
        suffix = "个" if self._t("lang_btn") == "EN" else "files"
        self.cnt.config(text=str(len(self.files)) + " " + suffix)


# ==============================================================================
# 主界面（纯 pack 布局）
# ==============================================================================

class App:

    def __init__(self, root):
        self.root     = root
        self.root.title("文档脱敏")
        self.root.configure(bg=BG)
        self.root.minsize(640, 600)
        self.scrubber = Scrubber()
        self.lang     = tk.StringVar(value="zh")   # "zh" | "en"
        self._ui_refs = {}   # 存储需要动态更新的控件引用
        self._build()
        self._check_deps()

    def _t(self, key, **kw):
        s = LANG[self.lang.get()][key]
        return s.format(**kw) if kw else s

    # ── 语言切换 ──
    def _toggle_lang(self):
        self.lang.set("en" if self.lang.get() == "zh" else "zh")
        self._refresh_lang()

    def _refresh_lang(self):
        r = self._ui_refs
        self.root.title(self._t("title"))
        r["lbl_title"].config(text=self._t("title"))
        r["lbl_subtitle"].config(text=self._t("subtitle"))
        r["lbl_sensitive"].config(text=self._t("sensitive_title"))
        r["lbl_sensitive_hint"].config(text=self._t("sensitive_hint"))
        r["lang_btn"].config(text=self._t("lang_btn"))
        r["chk_del"].config(text=self._t("del_json_opt"))
        r["status_lbl"].config(text=self._t("ready"))
        r["footer_lbl"].config(text=self._t("footer"))
        # 更新占位符
        cur = self.word_box.get("1.0", tk.END).strip()
        old_ph = LANG["en"]["sensitive_ph"] if self.lang.get() == "zh" else LANG["zh"]["sensitive_ph"]
        if cur == old_ph or cur == "":
            self.word_box.delete("1.0", tk.END)
            self.word_box.insert("1.0", self._t("sensitive_ph"))
            self.word_box.config(fg=T2)
            self._ph = self._t("sensitive_ph")
        else:
            self._ph = self._t("sensitive_ph")
        # 更新区域标签
        self._update_section_label(r["enc_zone_dot"], r["enc_zone_lbl"], "enc_zone")
        self._update_section_label(r["dec_zone_dot"], r["dec_zone_lbl"], "dec_zone")
        # 更新 FileList 按钮
        self.enc.lang_var = self.lang
        self.dec.lang_var = self.lang
        self.enc.update_lang()
        self.dec.update_lang()
        # 更新操作按钮（Canvas 需重建）
        self._rebuild_action_btns()
        # 更新 badge
        self._upd_badge()

    def _update_section_label(self, dot_canvas, lbl, key):
        lbl.config(text=self._t(key))

    def _rebuild_action_btns(self):
        r = self._ui_refs
        # 加密按钮
        for w in r["enc_btn_row"].winfo_children(): w.destroy()
        self._rndbtn(r["enc_btn_row"], self._t("start_enc"), ORANGE,
                     lambda: self._run("enc")).pack(side="left", padx=(0, 6))
        # 解密按钮
        for w in r["dec_btn_row"].winfo_children(): w.destroy()
        self._rndbtn(r["dec_btn_row"], self._t("start_dec"), BLUE,
                     lambda: self._run("dec")).pack(side="left", padx=(0, 6))
        self._rndbtn(r["dec_btn_row"], self._t("clear_map"), "#AEAEB2",
                     self._clear_map).pack(side="left")
        # 复制AI提示词按钮
        self._rebuild_prompt_btn()

    def _rebuild_prompt_btn(self):
        r = self._ui_refs
        frame = r["prompt_btn_frame"]
        for w in frame.winfo_children(): w.destroy()
        self._rndbtn_fullwidth(frame, self._t("copy_prompt"), PURPLE, self._copy_prompt)

    # ── 依赖检查 ──
    def _check_deps(self):
        miss = []
        if not DOCX_OK: miss.append("python-docx")
        if not PPTX_OK: miss.append("python-pptx")
        if not XLSX_OK: miss.append("openpyxl")
        if not DND_OK:  miss.append("tkinterdnd2" + ("（拖拽）" if self.lang.get()=="zh" else " (drag & drop)"))
        if miss:
            messagebox.showwarning(
                self._t("dep_title"),
                self._t("dep_body", libs="\n".join("  · " + m for m in miss))
            )

    # ── 构建界面（全部 pack）──
    def _build(self):
        r = self._ui_refs
        wrap = tk.Frame(self.root, bg=BG)
        wrap.pack(fill="both", expand=True, padx=20, pady=20)

        # ── 标题行 ──
        row0 = tk.Frame(wrap, bg=BG)
        row0.pack(fill="x", pady=(0, 14))

        r["lbl_title"] = tk.Label(row0, text=self._t("title"),
                                  font=(FONT, 18, "bold"), bg=BG, fg=T1)
        r["lbl_title"].pack(side="left")

        r["lbl_subtitle"] = tk.Label(row0, text=self._t("subtitle"),
                                     font=(FONT, 9), bg=BG, fg=T2)
        r["lbl_subtitle"].pack(side="left", padx=12, pady=4)

        # 右上角：语言切换按钮 + badge
        r["lang_btn"] = tk.Button(
            row0, text=self._t("lang_btn"),
            font=(FONT, 8, "bold"),
            bg=BLUE, fg="white",
            activebackground=FileList._tint(BLUE, 0.82),
            activeforeground="white",
            relief="flat", bd=0, padx=10, pady=3,
            cursor="hand2",
            command=self._toggle_lang
        )
        r["lang_btn"].pack(side="right", padx=(6, 0))

        self.badge = tk.Label(row0, text=self._t("badge_empty"),
                              font=(FONT, 8), bg="#E5E5EA", fg=T2,
                              padx=8, pady=2)
        self.badge.pack(side="right")

        # ── 敏感词卡片 ──
        wc = tk.Frame(wrap, bg=CARD, highlightbackground=BORDER,
                      highlightthickness=1)
        wc.pack(fill="x", pady=(0, 12))

        wi = tk.Frame(wc, bg=CARD)
        wi.pack(fill="x", padx=14, pady=10)

        r["lbl_sensitive"] = tk.Label(wi, text=self._t("sensitive_title"),
                                      font=(FONT, 10, "bold"), bg=CARD, fg=T1)
        r["lbl_sensitive"].pack(anchor="w")

        r["lbl_sensitive_hint"] = tk.Label(wi, text=self._t("sensitive_hint"),
                                           font=(FONT, 8), bg=CARD, fg=T2)
        r["lbl_sensitive_hint"].pack(anchor="w", pady=(1, 5))

        self.word_box = tk.Text(
            wi, height=2, font=(FONT, 10),
            relief="flat", bd=0,
            bg="#F2F2F7", fg=T1,
            insertbackground=ORANGE,
            wrap=tk.WORD,
            highlightbackground=BORDER,
            highlightthickness=1
        )
        self.word_box.pack(fill="x", ipady=5)
        self._ph = self._t("sensitive_ph")
        self.word_box.insert("1.0", self._ph)
        self.word_box.config(fg=T2)
        self.word_box.bind("<FocusIn>",  self._wi)
        self.word_box.bind("<FocusOut>", self._wo)

        # ── 双列文件区：使用 grid 确保等宽 ──
        cols = tk.Frame(wrap, bg=BG)
        cols.pack(fill="both", expand=True, pady=(0, 12))
        cols.columnconfigure(0, weight=1, uniform="col")
        cols.columnconfigure(1, weight=1, uniform="col")
        cols.rowconfigure(0, weight=1)

        # 左列：加密区
        left = tk.Frame(cols, bg=BG)
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 6))

        enc_lbl_frame = tk.Frame(left, bg=BG)
        enc_lbl_frame.pack(fill="x", pady=(0, 5))
        enc_dot = tk.Canvas(enc_lbl_frame, width=8, height=8, bg=BG, highlightthickness=0)
        enc_dot.pack(side="left", pady=1)
        enc_dot.create_oval(1, 1, 7, 7, fill=ORANGE, outline="")
        r["enc_zone_dot"] = enc_dot
        r["enc_zone_lbl"] = tk.Label(enc_lbl_frame, text=self._t("enc_zone"),
                                     font=(FONT, 10, "bold"), bg=BG, fg=T1)
        r["enc_zone_lbl"].pack(side="left", padx=5)

        enc_card = tk.Frame(left, bg=CARD, highlightbackground=BORDER, highlightthickness=1)
        enc_card.pack(fill="both", expand=True)

        self.enc = FileList(enc_card, ORANGE, "enc", self.scrubber, lang_var=self.lang)
        self.enc.pack(fill="both", expand=True, padx=8, pady=(8, 4))

        r["enc_btn_row"] = tk.Frame(enc_card, bg=CARD)
        r["enc_btn_row"].pack(fill="x", padx=8, pady=(0, 8))
        self._rndbtn(r["enc_btn_row"], self._t("start_enc"), ORANGE,
                     lambda: self._run("enc")).pack(side="left", padx=(0, 6))

        # 右列：解密区
        right = tk.Frame(cols, bg=BG)
        right.grid(row=0, column=1, sticky="nsew", padx=(6, 0))

        dec_lbl_frame = tk.Frame(right, bg=BG)
        dec_lbl_frame.pack(fill="x", pady=(0, 5))
        dec_dot = tk.Canvas(dec_lbl_frame, width=8, height=8, bg=BG, highlightthickness=0)
        dec_dot.pack(side="left", pady=1)
        dec_dot.create_oval(1, 1, 7, 7, fill=BLUE, outline="")
        r["dec_zone_dot"] = dec_dot
        r["dec_zone_lbl"] = tk.Label(dec_lbl_frame, text=self._t("dec_zone"),
                                     font=(FONT, 10, "bold"), bg=BG, fg=T1)
        r["dec_zone_lbl"].pack(side="left", padx=5)

        dec_card = tk.Frame(right, bg=CARD, highlightbackground=BORDER, highlightthickness=1)
        dec_card.pack(fill="both", expand=True)

        self.dec = FileList(dec_card, BLUE, "dec", self.scrubber,
                            on_json=self._json_loaded, lang_var=self.lang)
        self.dec.pack(fill="both", expand=True, padx=8, pady=(8, 4))

        r["dec_btn_row"] = tk.Frame(dec_card, bg=CARD)
        r["dec_btn_row"].pack(fill="x", padx=8, pady=(0, 8))
        self._rndbtn(r["dec_btn_row"], self._t("start_dec"), BLUE,
                     lambda: self._run("dec")).pack(side="left", padx=(0, 6))
        self._rndbtn(r["dec_btn_row"], self._t("clear_map"), "#AEAEB2",
                     self._clear_map).pack(side="left")

        # ── 选项行 ──
        opt_row = tk.Frame(wrap, bg=BG)
        opt_row.pack(fill="x", pady=(0, 10))

        self.del_json_var = tk.BooleanVar(value=False)
        r["chk_del"] = tk.Checkbutton(
            opt_row,
            text=self._t("del_json_opt"),
            variable=self.del_json_var,
            font=(FONT, 9), bg=BG, fg=T2,
            activebackground=BG,
            selectcolor=CARD,
            cursor="hand2"
        )
        r["chk_del"].pack(side="left")

        # ── 复制AI提示词按钮（全宽）──
        r["prompt_btn_frame"] = tk.Frame(wrap, bg=BG)
        r["prompt_btn_frame"].pack(fill="x", pady=(0, 10))
        self._rndbtn_fullwidth(r["prompt_btn_frame"], self._t("copy_prompt"), PURPLE, self._copy_prompt)

        # ── 状态栏 ──
        status_bar = tk.Frame(wrap, bg=CARD, highlightbackground=BORDER,
                              highlightthickness=1)
        status_bar.pack(fill="x")

        si = tk.Frame(status_bar, bg=CARD)
        si.pack(fill="x", padx=12, pady=7)

        r["status_lbl"] = tk.Label(si, text=self._t("ready"),
                                   font=(FONT, 9), bg=CARD, fg=T2)
        r["status_lbl"].pack(side="right", padx=(10, 0))
        self.status = r["status_lbl"]

        self.pv = tk.DoubleVar()
        style = ttk.Style()
        style.theme_use("default")
        style.configure("S.Horizontal.TProgressbar",
                        troughcolor="#E5E5EA",
                        background=ORANGE,
                        thickness=6, borderwidth=0)
        ttk.Progressbar(si, variable=self.pv, maximum=100,
                        mode="determinate",
                        style="S.Horizontal.TProgressbar").pack(side="left", fill="x", expand=True)

        # ── 底部提示 ──
        r["footer_lbl"] = tk.Label(self.root, text=self._t("footer"),
                                   fg="red", font=("微软雅黑", 9))
        r["footer_lbl"].pack(side="bottom", pady=5)

    def _rndbtn_fullwidth(self, parent, text, color, cmd, pady=9, radius=10):
        """全宽圆角按钮，随 parent 宽度自适应"""
        import tkinter.font as tkfont
        font_obj = (FONT, 10, "bold")
        f = tkfont.Font(family=FONT, size=10, weight="bold")
        th = f.metrics("linespace")
        h = th + pady * 2
        hover_color = FileList._tint(color, 0.82)

        cv = tk.Canvas(parent, height=h, bg=parent["bg"],
                       highlightthickness=0, cursor="hand2")
        cv.pack(fill="x", expand=True)

        current_color = [color]  # 用列表存储当前颜色，方便闭包修改

        def _draw(col=None, event=None):
            if col is None:
                col = current_color[0]
            else:
                current_color[0] = col
            # 优先用事件宽度，否则用 winfo_width
            if event is not None and hasattr(event, 'width'):
                w = event.width
            else:
                w = cv.winfo_width()
            if w < 10:
                return
            r = radius
            cv.delete("all")
            cv.create_arc(0, 0, 2*r, 2*r, start=90, extent=90, fill=col, outline=col)
            cv.create_arc(w-2*r, 0, w, 2*r, start=0, extent=90, fill=col, outline=col)
            cv.create_arc(0, h-2*r, 2*r, h, start=180, extent=90, fill=col, outline=col)
            cv.create_arc(w-2*r, h-2*r, w, h, start=270, extent=90, fill=col, outline=col)
            cv.create_rectangle(r, 0, w-r, h, fill=col, outline=col)
            cv.create_rectangle(0, r, w, h-r, fill=col, outline=col)
            cv.create_text(w//2, h//2, text=text, fill="white", font=font_obj)

        cv.bind("<Configure>", lambda e: _draw(event=e))
        cv.bind("<Button-1>",  lambda e: cmd())
        cv.bind("<Enter>",     lambda e: _draw(col=hover_color))
        cv.bind("<Leave>",     lambda e: _draw(col=color))

        # 强制延迟重绘：等布局完成后再画一次，解决初次/重建后不显示的问题
        def _force_redraw():
            _draw(col=color)
            # 若宽度还没更新，再等一帧
            if cv.winfo_width() < 10:
                cv.after(16, _force_redraw)

        cv.after(10, _force_redraw)

    def _rndbtn(self, parent, text, color, cmd, padx=14, pady=7, radius=10):
        """固定宽度圆角按钮"""
        font_obj = (FONT, 9, "bold")
        import tkinter.font as tkfont
        f = tkfont.Font(family=FONT, size=9, weight="bold")
        tw = f.measure(text)
        th = f.metrics("linespace")
        w = tw + padx * 2
        h = th + pady * 2
        hover_color = FileList._tint(color, 0.82)

        cv = tk.Canvas(parent, width=w, height=h,
                       bg=parent["bg"], highlightthickness=0, cursor="hand2")

        def _draw(col):
            cv.delete("all")
            r = radius
            cv.create_arc(0, 0, 2*r, 2*r, start=90, extent=90, fill=col, outline=col)
            cv.create_arc(w-2*r, 0, w, 2*r, start=0, extent=90, fill=col, outline=col)
            cv.create_arc(0, h-2*r, 2*r, h, start=180, extent=90, fill=col, outline=col)
            cv.create_arc(w-2*r, h-2*r, w, h, start=270, extent=90, fill=col, outline=col)
            cv.create_rectangle(r, 0, w-r, h, fill=col, outline=col)
            cv.create_rectangle(0, r, w, h-r, fill=col, outline=col)
            cv.create_text(w//2, h//2, text=text, fill="white", font=font_obj)

        _draw(color)
        cv.bind("<Button-1>", lambda e: cmd())
        cv.bind("<Enter>",    lambda e: _draw(hover_color))
        cv.bind("<Leave>",    lambda e: _draw(color))
        return cv

    def _section_label(self, parent, text, color):
        f = tk.Frame(parent, bg=BG)
        f.pack(fill="x", pady=(0, 5))
        c = tk.Canvas(f, width=8, height=8, bg=BG, highlightthickness=0)
        c.pack(side="left", pady=1)
        c.create_oval(1, 1, 7, 7, fill=color, outline="")
        tk.Label(f, text=text, font=(FONT, 10, "bold"),
                 bg=BG, fg=T1).pack(side="left", padx=5)

    # ── 输入框占位符 ──
    def _wi(self, e):
        if self.word_box.get("1.0", tk.END).strip() == self._ph:
            self.word_box.delete("1.0", tk.END)
            self.word_box.config(fg=T1)

    def _wo(self, e):
        if not self.word_box.get("1.0", tk.END).strip():
            self.word_box.insert("1.0", self._ph)
            self.word_box.config(fg=T2)

    # ── JSON 载入回调 ──
    def _json_loaded(self, n, path):
        self._upd_badge()
        self._set_status(self._t("map_loaded", n=n), BLUE)
        messagebox.showinfo(self._t("map_load_ok_title"),
                            self._t("map_load_ok_body", n=n))

    # ── 获取敏感词 ──
    def _words(self):
        raw = self.word_box.get("1.0", tk.END).strip()
        if raw == self._ph or not raw:
            return []
        return [w.strip() for w in re.split(r"[,，\n]", raw) if w.strip()]

    # ── 执行 ──
    def _run(self, mode):
        zone  = self.enc if mode == "enc" else self.dec
        files = zone.files

        if not files:
            key = "warn_no_file_enc" if mode == "enc" else "warn_no_file_dec"
            messagebox.showwarning(self._t("tip"), self._t(key))
            return

        if mode == "enc":
            words = self._words()
            if not words:
                messagebox.showwarning(self._t("tip"), self._t("warn_no_words"))
                return
        else:
            words = []
            if not self.scrubber.reverse and files:
                ok, mp = self.scrubber.auto_load(files[0])
                if ok:
                    self._upd_badge()
                    self._set_status(self._t("auto_load_map", name=os.path.basename(mp)), BLUE)
                else:
                    ans = messagebox.askokcancel(
                        self._t("warn_no_map_title"),
                        self._t("warn_no_map_body")
                    )
                    if ans:
                        jp = filedialog.askopenfilename(
                            title=self._t("pick_map_title"),
                            filetypes=[("JSON", "*.json"), ("All Files", "*.*")]
                        )
                        if jp:
                            try:
                                self.scrubber.load(jp)
                                self._upd_badge()
                            except Exception as ex:
                                messagebox.showerror(self._t("error"), str(ex))
                                return
                        else:
                            return
                    else:
                        return

        self.pv.set(0)
        self._set_status(self._t("processing"), T2)

        def worker():
            def cb(cur, tot, name, st):
                self.pv.set(cur / tot * 100)
                self._set_status(
                    "(" + str(cur) + "/" + str(tot) + ") " + name + " " + st, T2)

            results = self.scrubber.batch(files, words, mode, cb)
            self.root.after(0, lambda: self._done(results, mode))

        threading.Thread(target=worker, daemon=True).start()

    def _done(self, results, mode):
        ok   = [r for r in results if r[2] is None]
        fail = [r for r in results if r[2] is not None]
        color  = ORANGE if mode == "enc" else BLUE
        action = self._t("enc_done") if mode == "enc" else self._t("dec_done")

        self.pv.set(100)
        self._set_status(action + " " + self._t("success") + "：" + str(len(ok)) + "，" +
                         self._t("fail") + "：" + str(len(fail)), color)
        self._upd_badge()

        if mode == "dec" and ok and self.del_json_var.get():
            deleted = []
            for path in list(self.scrubber.json_paths):
                try:
                    if os.path.exists(path):
                        os.remove(path)
                        deleted.append(os.path.basename(path))
                except Exception:
                    pass
            if deleted:
                self._set_status(self._t("map_deleted", names=", ".join(deleted)), RED)

        lines = []
        if ok:
            lines.append(self._t("ok_n", n=len(ok)))
            for inp, out, _, mp in ok:
                lines.append("  " + os.path.basename(inp) +
                              "  →  " + os.path.basename(out))
        if fail:
            lines.append("\n" + self._t("fail_n", n=len(fail)))
            for inp, _, err, _ in fail:
                lines.append("  " + os.path.basename(inp) + "\n    " + str(err))

        if mode == "enc" and ok:
            lines.append(self._t("suggest_prompt"))

        title = action + (self._t("all_ok") if not fail else self._t("partial_fail"))
        if fail:
            messagebox.showwarning(title, "\n".join(lines))
        else:
            messagebox.showinfo(title, "\n".join(lines))

    # ── 复制AI提示词 ──
    def _copy_prompt(self):
        self.root.clipboard_clear()
        self.root.clipboard_append(AI_PROMPT)
        self.root.update()
        messagebox.showinfo(self._t("copied_title"), self._t("copied_body"))

    # ── 清空映射 ──
    def _clear_map(self):
        self.scrubber = Scrubber()
        self.enc.scrubber = self.scrubber
        self.dec.scrubber = self.scrubber
        self._upd_badge()
        self._set_status(self._t("map_cleared"), T2)
        messagebox.showinfo(self._t("map_clear_title"), self._t("map_clear_body"))

    # ── 辅助 ──
    def _set_status(self, text, color=None):
        self.status.config(text=text, fg=color or T2)

    def _upd_badge(self):
        n = len(self.scrubber.mapping)
        if n:
            self.badge.config(text=self._t("badge_n", n=n),
                              bg="#D1F0DA", fg=GREEN)
        else:
            self.badge.config(text=self._t("badge_empty"), bg="#E5E5EA", fg=T2)


# ==============================================================================
# 入口
# ==============================================================================

def main():
    root = TkinterDnD.Tk() if DND_OK else tk.Tk()
    App(root)
    root.mainloop()


if __name__ == "__main__":
    main()